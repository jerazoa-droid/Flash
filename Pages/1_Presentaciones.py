import os
import json
import re
import random
import requests
import streamlit as st
import anthropic
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image, ImageStat
from fpdf import FPDF

# Configuración
load_dotenv()
api_key = os.getenv("ANTHROPIC_API_KEY")
unsplash_key = os.getenv("UNSPLASH_ACCESS_KEY")

st.set_page_config(page_title="FlashClass - Presentaciones", layout="centered")

# Ocultar barra lateral
st.markdown("""
    <style>
        [data-testid="stSidebar"] { display: none; }
        [data-testid="collapsedControl"] { display: none; }
    </style>
""", unsafe_allow_html=True)

# --- FUNCIONES DE UTILIDAD ---
def sanear_texto(texto, max_palabra=45):
    if texto is None:
        return ""
    texto = texto.encode("latin-1", "replace").decode("latin-1")
    palabras = texto.split(" ")
    palabras_seguras = []
    for palabra in palabras:
        if len(palabra) > max_palabra:
            partes = [palabra[j:j + max_palabra] for j in range(0, len(palabra), max_palabra)]
            palabras_seguras.append(" ".join(partes))
        else:
            palabras_seguras.append(palabra)
    return " ".join(palabras_seguras)

def escribir_linea_pdf_segura(pdf, alto, texto, fallback_max=200):
    try:
        pdf.multi_cell(0, alto, text=texto)
    except Exception:
        texto_truncado = texto[:fallback_max] + ("..." if len(texto) > fallback_max else "")
        try:
            pdf.multi_cell(0, alto, text=texto_truncado)
        except Exception:
            pass

def crear_archivo_pdf_guia_apoyo(asignatura, tema, nivel):
    if not api_key: return None
    client = anthropic.Anthropic(api_key=api_key)
    prompt_guia_adicional = f"""
    Redacta una Guía de Apoyo Académico Complementaria, exhaustiva y de riguroso NIVEL UNIVERSITARIO sobre '{tema}' para la asignatura de '{asignatura}' (Contexto de dificultad base: {nivel}).
    
    REQUISITOS OBLIGATORIOS PARA LA GUÍA:
    1. Rigor conceptual, técnico y formal propio de educación superior.
    2. Si el tema involucra modelos matemáticos, físicos, químicos o de ingeniería, incluye explícitamente las ecuaciones analíticas completas, la deducción de expresiones clave, la definición detallada de cada variable, sus unidades y significado.
    3. Estructura el documento utilizando subtítulos claros con '##'.
    4. Incluye un apartado formal con un problema o ejercicio complejo resuelto paso a paso.
    """
    try:
        resp = client.messages.create(
            model="claude-sonnet-5", max_tokens=8192,
            messages=[{"role": "user", "content": prompt_guia_adicional}]
        )
        texto_guia = "".join(bloque.text for bloque in resp.content if hasattr(bloque, "text"))
    except Exception as e:
        texto_guia = f"Error al generar el contenido de la guía: {e}"

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("helvetica", "B", 16)
    pdf.cell(0, 10, text=sanear_texto(f"FlashClass - Guía de Apoyo: {tema}"), new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_font("helvetica", "I", 11)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 10, text=sanear_texto(f"Asignatura: {asignatura} | Nivel Universitario Riguroso"), new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(8)
    pdf.set_text_color(0, 0, 0)
    pdf.set_font("helvetica", "", 10)

    for linea in texto_guia.split("\n"):
        linea_limpia = linea.strip()
        if not linea_limpia:
            pdf.ln(3)
            continue
        if linea_limpia.startswith("## ") or linea_limpia.startswith("# "):
            pdf.set_font("helvetica", "B", 12)
            escribir_linea_pdf_segura(pdf, 7, sanear_texto(linea_limpia.replace("#", "").strip()))
            pdf.set_font("helvetica", "", 10)
            pdf.ln(2)
        else:
            escribir_linea_pdf_segura(pdf, 5.5, sanear_texto(linea_limpia))
    return bytes(pdf.output())

def es_imagen_oscura(imagen_bytes):
    try:
        img = Image.open(imagen_bytes).convert('L')
        estadisticas = ImageStat.Stat(img)
        return estadisticas.mean[0] < 127
    except Exception:
        return False

def crear_archivo_pptx(datos_json, archivo_plantilla_usuario=None, imagen_fondo_usuario=None):
    if archivo_plantilla_usuario is not None:
        prs = Presentation(BytesIO(archivo_plantilla_usuario.getvalue()))
    elif os.path.exists("plantilla_base.pptx"):
        prs = Presentation("plantilla_base.pptx")
    else:
        prs = Presentation()

    fondo_es_oscuro = False
    if imagen_fondo_usuario:
        fondo_es_oscuro = es_imagen_oscura(BytesIO(imagen_fondo_usuario.getvalue()))

    for diapo in datos_json["diapositivas"]:
        if diapo.get("mantener", True):
            layout_index = 6 if len(prs.slide_layouts) > 6 else 1
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            if imagen_fondo_usuario:
                try:
                    img_stream_fondo = BytesIO(imagen_fondo_usuario.getvalue())
                    slide.shapes.add_picture(img_stream_fondo, 0, 0, width=prs.slide_width, height=prs.slide_height)
                except Exception:
                    pass

            caja_titulo = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(1.0))
            tf_titulo = caja_titulo.text_frame
            tf_titulo.word_wrap = True
            p_titulo = tf_titulo.paragraphs[0]
            p_titulo.text = diapo["titulo"]
            p_titulo.font.size = Pt(28)
            p_titulo.font.bold = True
            if fondo_es_oscuro: p_titulo.font.color.rgb = RGBColor(255, 255, 255)

            tiene_imagen = bool(diapo.get("imagen_seleccionada") or diapo.get("archivo_local"))
            ancho_texto = Inches(5.0) if tiene_imagen else Inches(10.5)

            caja_contenido = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), ancho_texto, Inches(4.5))
            tf_contenido = caja_contenido.text_frame
            tf_contenido.word_wrap = True
            primera_linea = True
            for linea in diapo["contenido"].split("\n"):
                limpia = linea.strip(" -*")
                if limpia:
                    p = tf_contenido.paragraphs[0] if primera_linea else tf_contenido.add_paragraph()
                    primera_linea = False
                    p.text = "- " + limpia
                    p.font.size = Pt(18)
                    p.space_after = Pt(10)
                    if fondo_es_oscuro: p.font.color.rgb = RGBColor(255, 255, 255)

            if tiene_imagen:
                try:
                    if diapo.get("archivo_local"):
                        slide.shapes.add_picture(BytesIO(diapo["archivo_local"].getvalue()), Inches(6.2), Inches(2.0), width=Inches(5.8))
                    elif diapo.get("imagen_seleccionada"):
                        img_response = requests.get(diapo["imagen_seleccionada"], timeout=10)
                        if img_response.status_code == 200:
                            slide.shapes.add_picture(BytesIO(img_response.content), Inches(6.2), Inches(2.0), width=Inches(5.8))
                except Exception:
                    pass
    archivo_en_memoria = BytesIO()
    prs.save(archivo_en_memoria)
    archivo_en_memoria.seek(0)
    return archivo_en_memoria

def extraer_json_de_texto(texto):
    match = re.search(r'\{.*\}', texto, re.DOTALL)
    if match: return json.loads(match.group(0))
    return json.loads(texto)

def buscar_imagenes_unsplash(query):
    if not unsplash_key: return []
    try:
        url = f"https://api.unsplash.com/search/photos?query={query}&client_id={unsplash_key}&per_page=3&page={random.randint(1, 5)}&orientation=landscape"
        return [img["urls"]["regular"] for img in requests.get(url, timeout=10).json().get("results", [])]
    except Exception: return []

REGLA_EJERCICIOS = """
REGLA GENERAL SOBRE EJERCICIOS: Cada vez que generes una diapositiva con un ejercicio, la diapositiva INMEDIATAMENTE SIGUIENTE debe contener el desarrollo paso a paso y la respuesta final.
"""

def obtener_definicion_nivel(nivel, tema, carrera):
    definiciones = {
        "Introductorio": "Asume que el alumno NO tiene conocimiento previo del tema. Usa lenguaje simple y analogias.",
        "Intermedio": f"Enfoque en la APLICACION practica de '{tema}'. Incluye obligatoriamente un ejercicio practico tipico de la asignatura con su solucion paso a paso.",
        "Avanzado": f"Asume dominio previo. Analiza si '{tema}' en '{carrera}' es cuantitativa o conceptual, y genera problemas de diseno o calculo avanzado según corresponda."
    }
    return definiciones.get(nivel, "")

def generar_clase():
    if not api_key:
        st.error("Error: API Key no configurada.")
        return False
    client = anthropic.Anthropic(api_key=api_key)
    reglas_extra = "- Incluye una diapositiva final con 2 ejercicios de desarrollo.\n" if st.session_state.incluir_ejercicios else ""
    definicion_nivel = obtener_definicion_nivel(st.session_state.nivel, st.session_state.tema, st.session_state.carrera)

    prompt_docente = f"""
    Crea una presentacion sobre '{st.session_state.tema}' para '{st.session_state.carrera}' (Nivel: '{st.session_state.nivel}').
    {reglas_extra}
    INSTRUCCIONES ESPECIFICAS PARA EL NIVEL '{st.session_state.nivel}':
    {definicion_nivel}
    {REGLA_EJERCICIOS}
    REGLA ESTRICTA: Tu respuesta debe iniciar y terminar UNICAMENTE con un objeto JSON valido. NO agregues introducciones. Estructura exacta:
    {{
        "diapositivas": [
            {{"id": 1, "titulo": "Titulo", "contenido": "Punto 1\\nPunto 2", "sugiere_imagen": true, "keyword_imagen": "keyword in english"}}
        ]
    }}
    """
    NIVELES_EFFORT = {"Introductorio": "high", "Intermedio": "high", "Avanzado": "xhigh"}
    with st.spinner("FlashClass diseñando la clase (JSON)..."):
        try:
            respuesta = client.messages.create(
                model="claude-sonnet-5", max_tokens=8192,
                output_config={"effort": NIVELES_EFFORT.get(st.session_state.nivel, "high")},
                messages=[{"role": "user", "content": prompt_docente}]
            )
            if respuesta.stop_reason == "max_tokens":
                st.error("La clase se corto por limite de tokens.")
                return False
            texto_crudo = "".join(b.text for b in respuesta.content if hasattr(b, "text"))
            datos_json = extraer_json_de_texto(texto_crudo)
            for d in datos_json["diapositivas"]:
                d.update({"mantener": True, "imagen_seleccionada": None, "archivo_local": None, "opciones_imagenes": []})
            st.session_state.clase_generada = datos_json
            st.session_state.archivo_ppt_final = None
            st.session_state.archivo_pdf_guia = None
            return True
        except Exception as e:
            st.error(f"Error procesando la IA: {e}")
            return False

# --- ESTADO DE SESIÓN ---
if "clase_generada" not in st.session_state: st.session_state.clase_generada = None
if "carrera" not in st.session_state: st.session_state.carrera = "Ingenieria Civil"
if "tema" not in st.session_state: st.session_state.tema = "Leyes de la Termodinamica"
if "nivel" not in st.session_state: st.session_state.nivel = "Introductorio"
if "incluir_ejercicios" not in st.session_state: st.session_state.incluir_ejercicios = True
if "archivo_ppt_final" not in st.session_state: st.session_state.archivo_ppt_final = None
if "archivo_pdf_guia" not in st.session_state: st.session_state.archivo_pdf_guia = None

# --- BOTÓN VOLVER GENERAL ---
if st.button("⬅️ Volver al Inicio"):
    st.session_state.clase_generada = None
    st.session_state.archivo_ppt_final = None
    st.session_state.archivo_pdf_guia = None
    st.switch_page("app.py")

# --- INTERFAZ ---
if st.session_state.clase_generada is None:
    st.title("⚡ FlashClass - Presentaciones PPTX")
    with st.form("formulario_inicial"):
        st.session_state.carrera = st.text_input("Asignatura", value=st.session_state.carrera)
        st.session_state.tema = st.text_input("Tema de la clase", value=st.session_state.tema)
        st.session_state.nivel = st.selectbox("Nivel academico", ["Introductorio", "Intermedio", "Avanzado"], index=["Introductorio", "Intermedio", "Avanzado"].index(st.session_state.nivel))
        st.session_state.incluir_ejercicios = st.checkbox("Incluir ejercicios", value=st.session_state.incluir_ejercicios)
        if st.form_submit_button("Generar Borrador PPT", use_container_width=True):
            if generar_clase(): st.rerun()
else:
    st.title("⚡ FlashClass - Editor de Clase")

    with st.expander("⚙️ Ajustes, Parámetros y Diseño de Plantilla", expanded=False):
        c_param, c_diseno = st.columns(2)
        with c_param:
            st.markdown("**Ajustar Parámetros**")
            with st.form("formulario_lateral"):
                lat_carrera = st.text_input("Asignatura", value=st.session_state.carrera)
                lat_tema = st.text_input("Tema", value=st.session_state.tema)
                lat_nivel = st.selectbox("Nivel", ["Introductorio", "Intermedio", "Avanzado"], index=["Introductorio", "Intermedio", "Avanzado"].index(st.session_state.nivel))
                lat_ejercicios = st.checkbox("Ejercicios", value=st.session_state.incluir_ejercicios)
                if st.form_submit_button("Regenerar PPT"):
                    st.session_state.carrera = lat_carrera
                    st.session_state.tema = lat_tema
                    st.session_state.nivel = lat_nivel
                    st.session_state.incluir_ejercicios = lat_ejercicios
                    if generar_clase(): st.rerun()
        with c_diseno:
            st.markdown("**Diseño y Plantilla**")
            plantilla_subida = st.file_uploader("1. Cargar plantilla (.pptx)", type=["pptx"], key="plantilla_pptx")
            fondo_subido = st.file_uploader("2. O cargar fondo (PNG, JPG)", type=["png", "jpg", "jpeg"], key="fondo_img")

    datos = st.session_state.clase_generada
    for i, diapo in enumerate(datos["diapositivas"]):
        with st.expander(f"Diapositiva {i + 1}: {diapo['titulo']}", expanded=True):
            col_texto, col_img = st.columns([2, 2])
            with col_texto:
                diapo["titulo"] = st.text_input("Titulo", value=diapo["titulo"], key=f"t_{i}")
                diapo["contenido"] = st.text_area("Contenido", value=diapo["contenido"], key=f"c_{i}", height=150)
                diapo["mantener"] = st.checkbox("Incluir en el PPT final", value=diapo["mantener"], key=f"m_{i}")
            with col_img:
                tipo_fuente = st.radio("Fuente de imagen:", ["Buscar en Unsplash", "Subir"], key=f"fuente_{i}", horizontal=True)
                if "Unsplash" in tipo_fuente:
                    term = st.text_input("Palabra clave", value=diapo.get("keyword_imagen", "") if diapo.get("sugiere_imagen") else "", placeholder="Ej: physics", key=f"kb_{i}")
                    if term and st.button("Buscar", key=f"btn_{i}"):
                        diapo["opciones_imagenes"] = buscar_imagenes_unsplash(term)
                        diapo["archivo_local"] = None
                    if diapo.get("opciones_imagenes") and not diapo.get("imagen_seleccionada") and not diapo.get("archivo_local"):
                        c_f = st.columns(3)
                        for idx, url in enumerate(diapo["opciones_imagenes"]):
                            with c_f[idx]:
                                st.image(url, use_container_width=True)
                                if st.button("Elegir", key=f"sel_{i}_{idx}"):
                                    diapo["imagen_seleccionada"] = url
                                    st.rerun()
                else:
                    up_img = st.file_uploader("Sube PNG/JPG", type=["png", "jpg", "jpeg"], key=f"sub_{i}")
                    if up_img:
                        diapo["archivo_local"] = up_img
                        diapo["imagen_seleccionada"] = None
                if diapo.get("imagen_seleccionada"):
                    st.success("Imagen web acoplada")
                    st.image(diapo["imagen_seleccionada"], width=180)
                    if st.button("Quitar", key=f"del_web_{i}"):
                        diapo["imagen_seleccionada"] = None
                        st.rerun()
                elif diapo.get("archivo_local"):
                    st.success("Imagen local acoplada")
                    st.image(diapo["archivo_local"], width=180)
                    if st.button("Quitar", key=f"del_loc_{i}"):
                        diapo["archivo_local"] = None
                        st.rerun()

    st.markdown("---")
    st.subheader("Finalizar y Descargar")

    if st.button("Generar Presentacion PPTX", use_container_width=True, type="primary"):
        with st.spinner("Ensamblando presentación PowerPoint..."):
            st.session_state.archivo_ppt_final = crear_archivo_pptx(datos, plantilla_subida, fondo_subido)
        st.success("¡Presentación lista para descargar!")

    if st.session_state.archivo_ppt_final is not None:
        st.markdown("---")
        col_d1, col_d2 = st.columns(2)
        with col_d1:
            st.download_button("Descargar Presentacion (.pptx)", data=st.session_state.archivo_ppt_final, file_name=f"FlashClass_{st.session_state.tema.replace(' ', '_')}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
        with col_d2:
            if st.button("Generar guia de apoyo (Rigor universitario)", use_container_width=True):
                with st.spinner("Claude redactando guia avanzada..."):
                    st.session_state.archivo_pdf_guia = crear_archivo_pdf_guia_apoyo(st.session_state.carrera, st.session_state.tema, st.session_state.nivel)
                st.rerun()

    if st.session_state.archivo_pdf_guia is not None:
        st.success("Guia de apoyo lista:")
        st.download_button("Descargar Guia de Apoyo (.pdf)", data=st.session_state.archivo_pdf_guia, file_name=f"FlashClass_Guia_Apoyo_{st.session_state.tema.replace(' ', '_')}.pdf", mime="application/pdf", use_container_width=True)